# marketing_ai/file_ingest.py
"""Extraction de texte depuis les fichiers téléversés par l'utilisateur.

Formats gérés : .txt .md .csv .tsv .pdf .docx
Pur pip (pypdf + python-docx), aucune dépendance système à installer.
"""
import io
import os

# Limite de caractères par document et au total, pour ne pas surcharger le prompt.
MAX_CHARS_PER_FILE = 6000
MAX_CHARS_TOTAL = 12000

SUPPORTED = (".txt", ".md", ".csv", ".tsv", ".pdf", ".docx", ".pptx")

# Images pour la vision
IMAGE_EXT = (".png", ".jpg", ".jpeg", ".webp", ".gif")
MAX_IMAGES = 4                       # limite le cout / la taille du prompt
MAX_IMAGE_BYTES = 5 * 1024 * 1024    # 5 Mo par image


def extract_text_from_bytes(name: str, data: bytes) -> str:
    """Renvoie le texte extrait d'un fichier (donné par son nom + contenu binaire)."""
    ext = os.path.splitext(name)[1].lower()
    try:
        if ext in (".txt", ".md", ".csv", ".tsv"):
            text = data.decode("utf-8", errors="replace")
        elif ext == ".pdf":
            text = _extract_pdf(data)
        elif ext == ".docx":
            text = _extract_docx(data)
        elif ext == ".pptx":
            text = _extract_pptx(data)
        else:
            return ""
    except Exception as e:
        return f"[Erreur de lecture de {name}: {e}]"
    return text.strip()[:MAX_CHARS_PER_FILE]


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _extract_docx(data: bytes) -> str:
    import docx
    doc = docx.Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        lines.append(txt)
    return "\n".join(lines)


# ============================================================
# === Récupération de contenu depuis une URL ===
# ============================================================
def _normalize_url(url: str) -> str:
    """Réécrit les liens Google Slides/Docs vers leur export téléchargeable."""
    import re
    m = re.search(r"docs\.google\.com/presentation/d/([\w-]+)", url)
    if m:
        return f"https://docs.google.com/presentation/d/{m.group(1)}/export/pdf"
    m = re.search(r"docs\.google\.com/document/d/([\w-]+)", url)
    if m:
        return f"https://docs.google.com/document/d/{m.group(1)}/export?format=txt"
    return url


def _html_to_text(html_text: str) -> str:
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(html_text, "html.parser")
    for tag in soup(["script", "style", "noscript", "header", "footer", "nav"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    return "\n".join(lines)


def fetch_url(url: str, timeout: int = 20) -> str:
    """Télécharge une URL et en extrait le texte.

    Gère : PDF, Word (.docx), PowerPoint (.pptx), pages HTML, et les liens
    Google Slides/Docs partagés (via leur export). Renvoie un texte tronqué.
    """
    import requests
    target = _normalize_url(url)
    try:
        r = requests.get(target, timeout=timeout,
                         headers={"User-Agent": "Mozilla/5.0"})
        r.raise_for_status()
    except Exception as e:
        return f"[Erreur lors de la récupération de l'URL : {e}]"

    ctype = (r.headers.get("content-type") or "").lower()
    low = target.lower()
    try:
        if "pdf" in ctype or low.endswith(".pdf"):
            text = _extract_pdf(r.content)
        elif "wordprocessing" in ctype or low.endswith(".docx"):
            text = _extract_docx(r.content)
        elif "presentation" in ctype or low.endswith(".pptx"):
            text = _extract_pptx(r.content)
        elif "html" in ctype or low.startswith("http"):
            text = _html_to_text(r.text)
        else:
            text = r.text
    except Exception as e:
        return f"[Erreur d'extraction du contenu de l'URL : {e}]"
    return text.strip()[:MAX_CHARS_TOTAL]


def _read_bytes(f):
    if hasattr(f, "getvalue"):
        return f.getvalue()
    if hasattr(f, "read"):
        return f.read()
    return bytes(f)


def images_to_data_uris(files) -> list:
    """Convertit des images téléversées en data-URI base64 pour la vision.

    Renvoie une liste de chaînes 'data:image/...;base64,...' (max MAX_IMAGES).
    """
    import base64
    uris = []
    for f in files:
        if len(uris) >= MAX_IMAGES:
            break
        name = getattr(f, "name", "image")
        ext = os.path.splitext(name)[1].lower()
        if ext not in IMAGE_EXT:
            continue
        data = _read_bytes(f)
        if not data or len(data) > MAX_IMAGE_BYTES:
            continue
        mime = getattr(f, "type", None) or f"image/{ext.lstrip('.').replace('jpg', 'jpeg')}"
        uris.append(f"data:{mime};base64,{base64.b64encode(data).decode()}")
    return uris


def combine_files(files) -> str:
    """Concatène le texte de plusieurs fichiers en un seul bloc de contexte.

    `files` : itérable d'objets ayant .name et une méthode de lecture des octets
    (getvalue() pour Streamlit, ou read()). Renvoie un texte tronqué au global.
    """
    parts = []
    for f in files:
        name = getattr(f, "name", "document")
        if hasattr(f, "getvalue"):
            data = f.getvalue()
        elif hasattr(f, "read"):
            data = f.read()
        else:
            data = bytes(f)
        text = extract_text_from_bytes(name, data)
        if text:
            parts.append(f"--- {name} ---\n{text}")
    combined = "\n\n".join(parts)
    return combined[:MAX_CHARS_TOTAL]
